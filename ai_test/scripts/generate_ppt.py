"""
生成 AiProtect 平台周会汇报 PPT（通俗版）
用 STAR 法则讲清楚 AI 生成测试用例：场景-任务-行动-结果
面向小白，通俗易懂
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

PRIMARY = RGBColor(0x1A, 0x73, 0xE8)
SECONDARY = RGBColor(0x34, 0xA8, 0x53)
ACCENT = RGBColor(0xFB, 0xBC, 0x04)
DARK = RGBColor(0x20, 0x2A, 0x36)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
RED = RGBColor(0xEA, 0x43, 0x35)
PURPLE = RGBColor(0x9C, 0x27, 0xB0)
ORANGE = RGBColor(0xFF, 0x57, 0x22)
TEAL = RGBColor(0x00, 0x96, 0x88)
PINK = RGBColor(0xE9, 0x1E, 0x63)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = align
    return tf


def add_bullet(tf, text, font_size=16, color=DARK, bold=False, space_before=Pt(6), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.font.bold = bold
    p.space_before = space_before
    p.space_after = space_after
    return p


def first_or_add(tf, text, font_size=14, color=DARK, bold=False, space_before=Pt(6)):
    if tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.font.bold = bold
    else:
        add_bullet(tf, text, font_size=font_size, color=color, bold=bold, space_before=space_before)


def add_header_bar(slide, title_text):
    bar = add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), PRIMARY)
    set_text(bar, title_text, font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    bar.text_frame.paragraphs[0].space_before = Pt(10)
    bar.text_frame.margin_left = Inches(0.6)
    bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def make_card(slide, x, y, w, h, title, title_color, items, item_size=14, item_color=DARK):
    card = add_rounded_rect(slide, x, y, w, h, WHITE, RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(slide, x, y, w, Inches(0.06), title_color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.4),
                 title, font_size=17, bold=True, color=title_color)
    tf = add_text_box(slide, x + Inches(0.2), y + Inches(0.6), w - Inches(0.4), h - Inches(0.8),
                      "", font_size=item_size, color=item_color)
    for item in items:
        first_or_add(tf, item, font_size=item_size, color=item_color, space_before=Pt(5))
    return card


# ============================================================
# Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, PRIMARY)
add_rect(slide, Inches(0.15), Inches(2.3), Inches(8), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(2.6), Inches(10), Inches(1.2),
             "AiProtect — AI 驱动的智能测试平台", font_size=40, bold=True, color=DARK)
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(10), Inches(0.8),
             "AI 生成测试用例：从痛点到落地实践", font_size=24, bold=False, color=MEDIUM_GRAY)
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(10), Inches(0.5),
             "用 STAR 法则讲清楚：场景 → 任务 → 行动 → 结果", font_size=18, bold=False, color=PRIMARY)
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(10), Inches(0.5),
             "汇报人：阮龙博", font_size=18, bold=False, color=MEDIUM_GRAY)


# ============================================================
# Slide 2: 今天聊什么（目录页）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "今天聊什么？")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
             "围绕 AI 生成测试用例，用 STAR 法则讲透这件事", font_size=18, color=MEDIUM_GRAY)

agenda = [
    ("S — Situation  场景", "我们为什么要做这件事？测试工作有哪些痛点？", PRIMARY, "01"),
    ("T — Task  任务", "这个平台到底能干什么？核心功能是什么？", SECONDARY, "02"),
    ("A — Action  行动", "具体怎么用？手把手操作演示", ORANGE, "03"),
    ("A — Action  进阶", "如何用得更好？背后的技术原理", PURPLE, "04"),
    ("R — Result  结果", "目前的效果、痛点和未来规划", TEAL, "05"),
]

for i, (title, desc, color, num) in enumerate(agenda):
    y = Inches(2.1) + i * Inches(1.0)
    # 编号圆圈用方块代替
    num_box = add_rounded_rect(slide, Inches(0.8), y, Inches(0.7), Inches(0.7), color)
    num_box.line.fill.background()
    set_text(num_box, num, font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    num_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, Inches(1.8), y + Inches(0.0), Inches(4), Inches(0.4),
                 title, font_size=20, bold=True, color=color)
    add_text_box(slide, Inches(1.8), y + Inches(0.38), Inches(10), Inches(0.4),
                 desc, font_size=16, color=MEDIUM_GRAY)


# ============================================================
# Slide 3: S — Situation 痛点场景
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "S — Situation  |  测试工作的痛点场景")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
             "先说说我们测试同学日常遇到的那些头疼事...", font_size=18, color=MEDIUM_GRAY)

pain_points = [
    {
        "icon": "😫", "title": "写用例太耗时",
        "desc": "一个需求文档几十页，光写功能用例就要 1-2 天\n正向、边界、异常，一个都不能漏\n需求变了还得全部重写",
        "color": RED,
    },
    {
        "icon": "🤯", "title": "需求理解偏差",
        "desc": "PRD 写得模糊，不同人理解不一样\n评审会上说的补充内容容易遗忘\n跨团队信息同步不及时",
        "color": ORANGE,
    },
    {
        "icon": "😰", "title": "覆盖率难保证",
        "desc": "凭经验容易漏测边界和异常场景\n没有系统性的覆盖率检查机制\n关键功能漏测 → 线上 Bug → 背锅",
        "color": PURPLE,
    },
    {
        "icon": "📋", "title": "格式不统一",
        "desc": "每个人写用例风格不同\n用例粒度忽粗忽细\n新人上手慢，质量参差不齐",
        "color": TEAL,
    },
]

card_w = Inches(2.85)
card_h = Inches(3.5)
gap = Inches(0.3)
start_x = Inches(0.5)
y = Inches(2.1)

for i, pp in enumerate(pain_points):
    x = start_x + i * (card_w + gap)
    card = add_rounded_rect(slide, x, y, card_w, card_h, WHITE, RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(slide, x, y, card_w, Inches(0.06), pp["color"])

    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), card_w - Inches(0.3), Inches(0.5),
                 f'{pp["icon"]}  {pp["title"]}', font_size=18, bold=True, color=pp["color"])

    tf = add_text_box(slide, x + Inches(0.15), y + Inches(0.8), card_w - Inches(0.3), card_h - Inches(1.0),
                      "", font_size=14, color=MEDIUM_GRAY)
    for line in pp["desc"].split('\n'):
        first_or_add(tf, line, font_size=14, color=MEDIUM_GRAY, space_before=Pt(6))

# 底部总结
summary_box = add_rounded_rect(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
                                RGBColor(0xFD, 0xE8, 0xE8))
summary_box.line.fill.background()
add_text_box(slide, Inches(0.8), Inches(6.15), Inches(11.8), Inches(0.7),
             "核心矛盾：需求越来越多、迭代越来越快，但测试人力和时间没有增加 → 需要 AI 来帮忙",
             font_size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 4: T — Task 平台能干什么
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "T — Task  |  AiProtect 平台能帮你做什么？")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
             "一句话：让 AI 帮你写测试用例，你只需要给它需求文档", font_size=18, color=MEDIUM_GRAY)

# 核心能力
core_features = [
    ("AI 功能用例生成", "丢一份需求文档进去\nAI 自动提取测试点\n按场景分组生成完整用例", PRIMARY, "核心功能"),
    ("AI 接口用例生成", "导入 Swagger/OpenAPI\n自动生成可执行接口用例\n支持断言和数据库校验", SECONDARY, "核心功能"),
    ("智能知识库", "上传文档沉淀业务知识\nAI 生成时自动检索参考\n评审记录也能利用起来", PURPLE, "辅助增强"),
    ("需求智能分析", "PRD 结构化解析\n自动生成业务流程图\n让 AI 先「读懂」需求", ORANGE, "辅助增强"),
    ("接口自动化执行", "用例编排与一键执行\n环境变量管理\n执行报告与日志", TEAL, "自动化"),
    ("XMind / 报告导出", "用例导出为 XMind 脑图\n测试日报自动生成\n飞书同步通知", PINK, "效率工具"),
]

card_w = Inches(3.7)
card_h = Inches(2.3)
gap_x = Inches(0.4)
gap_y = Inches(0.35)
start_x = Inches(0.6)
start_y = Inches(2.0)

for i, (title, desc, color, tag) in enumerate(core_features):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    card = add_rounded_rect(slide, x, y, card_w, card_h, WHITE, RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(slide, x, y, card_w, Inches(0.06), color)

    # tag 标签
    tag_box = add_rounded_rect(slide, x + card_w - Inches(1.2), y + Inches(0.12), Inches(1.0), Inches(0.3), color)
    tag_box.line.fill.background()
    set_text(tag_box, tag, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tag_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(1.5), Inches(0.4),
                 title, font_size=17, bold=True, color=color)

    tf = add_text_box(slide, x + Inches(0.15), y + Inches(0.6), card_w - Inches(0.3), card_h - Inches(0.8),
                      "", font_size=13, color=MEDIUM_GRAY)
    for line in desc.split('\n'):
        first_or_add(tf, f"•  {line}", font_size=13, color=MEDIUM_GRAY, space_before=Pt(4))

# 技术栈
add_text_box(slide, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
             "技术栈：Vue 3 + Element Plus  |  FastAPI + LangChain + LangGraph  |  MySQL 8  |  Docker",
             font_size=13, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 5: A — Action 怎么用（操作流程）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "A — Action  |  具体怎么用？（3 步搞定）")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
             "不需要懂 AI，只需要 3 步，像聊天一样简单", font_size=18, color=MEDIUM_GRAY)

steps = [
    {
        "num": "Step 1",
        "title": "给 AI 喂需求",
        "color": PRIMARY,
        "items": [
            "上传需求文档（PDF/Word/Markdown）",
            "或者直接粘贴需求文字",
            "也可以传截图、设计稿、评审视频",
            "还可以勾选知识库里的相关文档",
            "",
            "支持 5 种输入方式：",
            "📄 文档  🖼️ 图片  🎬 视频  📝 文本  📚 知识库",
        ]
    },
    {
        "num": "Step 2",
        "title": "AI 提取测试点",
        "color": SECONDARY,
        "items": [
            "点击「AI 生成测试点」按钮",
            "AI 实时流式输出，边生成边看",
            "自动分为 3 类：",
            "  🟢 正向验证 — 功能正常场景",
            "  🟡 边界测试 — 极限/边界条件",
            "  🔴 异常处理 — 错误/异常场景",
            "",
            "可以人工审核和编辑测试点",
        ]
    },
    {
        "num": "Step 3",
        "title": "一键生成用例",
        "color": ORANGE,
        "items": [
            "点击「生成用例」按钮",
            "AI 自动按场景分组生成用例",
            "每条用例包含完整的八要素：",
            "  编号、名称、优先级、前置条件",
            "  测试步骤、输入数据、预期结果",
            "",
            "支持导出 XMind 脑图",
            "也可以直接在线管理和执行",
        ]
    },
]

card_w = Inches(3.7)
card_h = Inches(5.0)
gap = Inches(0.4)
start_x = Inches(0.6)
y = Inches(2.0)

for i, step in enumerate(steps):
    x = start_x + i * (card_w + gap)

    card = add_rounded_rect(slide, x, y, card_w, card_h, WHITE, RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(slide, x, y, card_w, Inches(0.06), step["color"])

    # 步骤编号
    num_box = add_rounded_rect(slide, x + Inches(0.15), y + Inches(0.15), Inches(1.0), Inches(0.35), step["color"])
    num_box.line.fill.background()
    set_text(num_box, step["num"], font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    num_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, x + Inches(1.3), y + Inches(0.15), card_w - Inches(1.5), Inches(0.4),
                 step["title"], font_size=18, bold=True, color=step["color"])

    tf = add_text_box(slide, x + Inches(0.2), y + Inches(0.65), card_w - Inches(0.4), card_h - Inches(0.9),
                      "", font_size=14, color=DARK)
    for item in step["items"]:
        first_or_add(tf, item, font_size=14, color=DARK if item else MEDIUM_GRAY, space_before=Pt(4))

# 箭头提示
add_text_box(slide, Inches(4.5), Inches(4.2), Inches(0.5), Inches(0.4),
             "→", font_size=30, bold=True, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.9), Inches(4.2), Inches(0.5), Inches(0.4),
             "→", font_size=30, bold=True, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 6: A — Action 进阶：背后的技术原理（通俗版）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "A — Action 进阶  |  AI 背后做了什么？（通俗版）")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
             "你点一个按钮，AI 在后台其实做了很多事，就像一个经验丰富的老测试在帮你干活",
             font_size=18, color=MEDIUM_GRAY)

# 左侧：流程
flow_box = add_rounded_rect(slide, Inches(0.5), Inches(2.0), Inches(6.0), Inches(5.2),
                             LIGHT_BLUE)
flow_box.line.fill.background()
add_text_box(slide, Inches(0.7), Inches(2.1), Inches(5.5), Inches(0.4),
             "AI 的工作流程（自动循环，直到覆盖全面）", font_size=18, bold=True, color=PRIMARY)

flow_steps = [
    ("1. 先做功课", "AI 先去知识库翻资料，看看评审会议记录，\n参考历史用例，「吃透业务」再动手"),
    ("2. 提取测试点", "像老测试一样，从需求里挖出所有需要测的点：\n正常场景、边界条件、异常情况一个不漏"),
    ("3. 自我检查", "AI 自己检查一遍：测试点覆盖全了吗？\n没覆盖全 → 自动补充 → 再检查（循环）"),
    ("4. 生成用例", "按业务场景分组，生成标准格式的测试用例\n每条用例都有完整的八要素"),
    ("5. 再次检查", "AI 再检查：每个测试点都有对应用例吗？\n没覆盖全 → 自动补充用例 → 再检查"),
    ("6. 保存入库", "全部覆盖 → 保存到数据库 → 完工！"),
]

y_pos = Inches(2.7)
for title, desc in flow_steps:
    add_text_box(slide, Inches(0.8), y_pos, Inches(5.4), Inches(0.28),
                 title, font_size=14, bold=True, color=DARK)
    add_text_box(slide, Inches(0.8), y_pos + Inches(0.28), Inches(5.4), Inches(0.45),
                 desc, font_size=12, color=MEDIUM_GRAY)
    y_pos += Inches(0.75)

# 右侧：知识增强
right_box = add_rounded_rect(slide, Inches(6.9), Inches(2.0), Inches(6.0), Inches(5.2),
                              RGBColor(0xE8, 0xF5, 0xE9))
right_box.line.fill.background()
add_text_box(slide, Inches(7.1), Inches(2.1), Inches(5.5), Inches(0.4),
             "知识增强 = 让 AI 不是「裸考」", font_size=18, bold=True, color=SECONDARY)

add_text_box(slide, Inches(7.1), Inches(2.7), Inches(5.5), Inches(0.5),
             "普通方式：只给 AI 一份需求文档 → AI 容易瞎编",
             font_size=14, color=RED)
add_text_box(slide, Inches(7.1), Inches(3.2), Inches(5.5), Inches(0.5),
             "我们的方式：给 AI 准备「4 份参考资料」再让它写",
             font_size=14, bold=True, color=SECONDARY)

sources = [
    ("📋  原始需求", "数据库中的结构化需求信息", PRIMARY),
    ("🔍  RAG 知识库", "自动检索相关的技术文档、业务文档", SECONDARY),
    ("💬  评审会议知识", "需求/技术/用例评审中的关键决策和遗漏", PURPLE),
    ("📚  历史用例", "同项目的历史用例，避免遗漏和重复", ORANGE),
]

y_pos = Inches(3.9)
for title, desc, color in sources:
    src_box = add_rounded_rect(slide, Inches(7.3), y_pos, Inches(5.3), Inches(0.7), WHITE, RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(slide, Inches(7.3), y_pos, Inches(0.08), Inches(0.7), color)
    add_text_box(slide, Inches(7.5), y_pos + Inches(0.05), Inches(5.0), Inches(0.3),
                 title, font_size=14, bold=True, color=color)
    add_text_box(slide, Inches(7.5), y_pos + Inches(0.35), Inches(5.0), Inches(0.3),
                 desc, font_size=12, color=MEDIUM_GRAY)
    y_pos += Inches(0.78)

add_text_box(slide, Inches(7.1), Inches(7.05), Inches(5.6), Inches(0.3),
             "→ 4 份资料合并成「增强需求文档」，再交给 AI 生成",
             font_size=13, bold=True, color=SECONDARY)


# ============================================================
# Slide 7: A — Action 进阶：如何用得更好
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "A — Action 进阶  |  如何让 AI 用例写得更好？")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
             "除了提示词调优，还有这些方法能大幅提升 AI 生成用例的质量",
             font_size=18, color=MEDIUM_GRAY)

better_methods = [
    {
        "title": "① 让 AI 吃透业务",
        "subtitle": "解决：AI 不懂业务、瞎编内容",
        "color": PRIMARY,
        "items": [
            "搭建测试专属 RAG 知识库 ✅ 已实现",
            "把需求文档、测试规范、历史案例都喂给 AI",
            "结合评审会议知识，避免遗漏 ✅ 已实现",
            "用 Schema 固定用例格式标准 ✅ 已实现",
        ]
    },
    {
        "title": "② 优化输入质量",
        "subtitle": "解决：需求理解偏差",
        "color": SECONDARY,
        "items": [
            "支持多模态输入（文档+图片+视频）✅ 已实现",
            "需求结构化解析，消除歧义 ✅ 已实现",
            "给不同功能标注风险等级和重点",
            "提供接口文档、原型图辅助理解",
        ]
    },
    {
        "title": "③ 输出质量检查",
        "subtitle": "解决：用例无效、覆盖不全",
        "color": ORANGE,
        "items": [
            "覆盖率双重验证（测试点+用例）✅ 已实现",
            "自动循环补充直到全覆盖 ✅ 已实现",
            "多模型交叉核对（规划中）",
            "测试环境实际运行验证（规划中）",
        ]
    },
    {
        "title": "④ 持续优化闭环",
        "subtitle": "解决：AI 不会越用越准",
        "color": PURPLE,
        "items": [
            "好用例 → 沉淀到知识库 → AI 学习参考",
            "坏用例 → 标注反馈 → 避免再犯",
            "打通需求/用例/缺陷管理工具",
            "用数据对比不同配置效果",
        ]
    },
]

card_w = Inches(5.9)
card_h = Inches(2.4)
positions = [
    (Inches(0.5), Inches(2.0)),
    (Inches(6.9), Inches(2.0)),
    (Inches(0.5), Inches(4.7)),
    (Inches(6.9), Inches(4.7)),
]

for i, method in enumerate(better_methods):
    x, y = positions[i]
    card = add_rounded_rect(slide, x, y, card_w, card_h, WHITE, RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(slide, x, y, card_w, Inches(0.06), method["color"])

    add_text_box(slide, x + Inches(0.2), y + Inches(0.12), card_w - Inches(0.4), Inches(0.35),
                 method["title"], font_size=17, bold=True, color=method["color"])
    add_text_box(slide, x + Inches(0.2), y + Inches(0.45), card_w - Inches(0.4), Inches(0.3),
                 method["subtitle"], font_size=12, color=MEDIUM_GRAY)

    tf = add_text_box(slide, x + Inches(0.2), y + Inches(0.8), card_w - Inches(0.4), card_h - Inches(1.0),
                      "", font_size=13, color=DARK)
    for item in method["items"]:
        first_or_add(tf, f"•  {item}", font_size=13, color=DARK, space_before=Pt(4))


# ============================================================
# Slide 8: R — Result 效果与对比
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "R — Result  |  效果怎么样？")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
             "AI 生成 vs 人工编写，对比看看差别",
             font_size=18, color=MEDIUM_GRAY)

# 对比表格用卡片模拟
# 人工
manual_box = add_rounded_rect(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.5),
                               RGBColor(0xFD, 0xE8, 0xE8))
manual_box.line.fill.background()
add_text_box(slide, Inches(0.7), Inches(2.1), Inches(5.3), Inches(0.4),
             "👤  传统人工编写", font_size=20, bold=True, color=RED)

manual_items = [
    ("⏱️  耗时", "一个中等需求 1-2 天", RED),
    ("📊  覆盖率", "依赖个人经验，容易漏测", RED),
    ("📋  格式", "因人而异，粒度不统一", RED),
    ("🔄  变更响应", "需求变更需大量重写", RED),
    ("📈  知识沉淀", "经验在个人脑中，难以传承", RED),
]

y_pos = Inches(2.7)
for icon_title, desc, color in manual_items:
    add_text_box(slide, Inches(0.8), y_pos, Inches(5.2), Inches(0.3),
                 icon_title, font_size=15, bold=True, color=DARK)
    add_text_box(slide, Inches(0.8), y_pos + Inches(0.3), Inches(5.2), Inches(0.3),
                 desc, font_size=14, color=color)
    y_pos += Inches(0.7)

# AI
ai_box = add_rounded_rect(slide, Inches(6.9), Inches(2.0), Inches(6.0), Inches(4.5),
                            RGBColor(0xE8, 0xF5, 0xE9))
ai_box.line.fill.background()
add_text_box(slide, Inches(7.1), Inches(2.1), Inches(5.5), Inches(0.4),
             "🤖  AI 辅助生成", font_size=20, bold=True, color=SECONDARY)

ai_items = [
    ("⏱️  耗时", "几分钟内完成，节省 80%+ 时间", SECONDARY),
    ("📊  覆盖率", "系统性覆盖：正向+边界+异常，双重验证", SECONDARY),
    ("📋  格式", "统一标准八要素，按场景分组", SECONDARY),
    ("🔄  变更响应", "重新输入需求，快速重新生成", SECONDARY),
    ("📈  知识沉淀", "知识库持续积累，AI 越用越准", SECONDARY),
]

y_pos = Inches(2.7)
for icon_title, desc, color in ai_items:
    add_text_box(slide, Inches(7.2), y_pos, Inches(5.4), Inches(0.3),
                 icon_title, font_size=15, bold=True, color=DARK)
    add_text_box(slide, Inches(7.2), y_pos + Inches(0.3), Inches(5.4), Inches(0.3),
                 desc, font_size=14, color=color)
    y_pos += Inches(0.7)

# 底部总结
add_text_box(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
             "💡 AI 不是要取代测试工程师，而是让测试工程师从「写用例」中解放出来，专注于更有价值的测试设计和分析",
             font_size=16, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 9: 目前的痛点与未来规划
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "R — Result  |  目前的痛点和未来规划")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
             "诚实地说说目前还有哪些不足，以及接下来要做什么",
             font_size=18, color=MEDIUM_GRAY)

# 当前痛点
pain_box = add_rounded_rect(slide, Inches(0.5), Inches(2.0), Inches(6.0), Inches(5.0),
                             RGBColor(0xFD, 0xF0, 0xE2))
pain_box.line.fill.background()
add_text_box(slide, Inches(0.7), Inches(2.1), Inches(5.5), Inches(0.4),
             "⚠️  目前存在的痛点", font_size=20, bold=True, color=ORANGE)

current_pains = [
    ("AI 偶尔会「幻觉」", "用例中可能出现不存在的功能或数据\n需要人工审核把关，不能完全依赖 AI"),
    ("复杂业务场景理解不够深", "涉及多系统交互、复杂链路的场景\nAI 可能只覆盖表面，深层逻辑靠人补充"),
    ("生成速度受限于模型", "大型需求文档生成耗时较长\n模型调用成本和响应时间需要平衡"),
    ("缺少反馈闭环", "用例好不好用，还没有系统化的评价机制\nAI 无法从使用反馈中自动学习改进"),
]

y_pos = Inches(2.7)
for title, desc in current_pains:
    add_text_box(slide, Inches(0.8), y_pos, Inches(5.4), Inches(0.28),
                 f"•  {title}", font_size=14, bold=True, color=DARK)
    add_text_box(slide, Inches(1.1), y_pos + Inches(0.3), Inches(5.1), Inches(0.5),
                 desc, font_size=12, color=MEDIUM_GRAY)
    y_pos += Inches(0.78)

# 未来规划
plan_box = add_rounded_rect(slide, Inches(6.9), Inches(2.0), Inches(6.0), Inches(5.0),
                             RGBColor(0xE8, 0xF5, 0xE9))
plan_box.line.fill.background()
add_text_box(slide, Inches(7.1), Inches(2.1), Inches(5.5), Inches(0.4),
             "🚀  未来规划", font_size=20, bold=True, color=SECONDARY)

plans = [
    {
        "phase": "短期（优先落地）",
        "color": SECONDARY,
        "items": [
            "完善用例反馈机制，好/坏用例标注",
            "丰富知识库内容，提升 RAG 检索质量",
            "优化 Prompt，减少幻觉和无效用例",
        ]
    },
    {
        "phase": "中期",
        "color": PRIMARY,
        "items": [
            "多模型交叉验证，提升用例准确性",
            "测试环境实际运行验证用例可执行性",
            "按场景自动选择最优模型",
        ]
    },
    {
        "phase": "长期",
        "color": PURPLE,
        "items": [
            "基于公司用例数据微调专属模型",
            "全链路自动化：需求→用例→执行→报告",
            "AI 主动发现潜在风险和漏测场景",
        ]
    },
]

y_pos = Inches(2.7)
for plan in plans:
    add_text_box(slide, Inches(7.2), y_pos, Inches(5.4), Inches(0.28),
                 plan["phase"], font_size=14, bold=True, color=plan["color"])
    y_pos += Inches(0.32)
    for item in plan["items"]:
        add_text_box(slide, Inches(7.4), y_pos, Inches(5.2), Inches(0.25),
                     f"•  {item}", font_size=12, color=DARK)
        y_pos += Inches(0.28)
    y_pos += Inches(0.15)


# ============================================================
# Slide 10: 落地建议总结
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_header_bar(slide, "落地优先级建议")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
             "不要想着一步到位，分阶段落地，先解决最痛的问题",
             font_size=18, color=MEDIUM_GRAY)

# 三个阶段
phases = [
    {
        "title": "🏃  立即可用（已实现）",
        "subtitle": "成本低、见效快",
        "color": SECONDARY,
        "bg": RGBColor(0xE8, 0xF5, 0xE9),
        "items": [
            "RAG 知识库 + 评审知识增强",
            "多模态输入（文档/图片/视频/文本）",
            "LangGraph 双层工作流 + 覆盖率验证",
            "按场景分组 + 标准八要素输出",
            "XMind 导出 + 在线用例管理",
        ]
    },
    {
        "title": "📋  近期优化（1-2 个月）",
        "subtitle": "提升质量和体验",
        "color": PRIMARY,
        "bg": LIGHT_BLUE,
        "items": [
            "用例质量反馈闭环（好/坏标注）",
            "知识库内容丰富和检索优化",
            "Prompt 持续调优，减少幻觉",
            "接口用例自动化执行完善",
            "测试报告和飞书通知优化",
        ]
    },
    {
        "title": "🚀  中长期规划（3-6 个月）",
        "subtitle": "深度智能化",
        "color": PURPLE,
        "bg": RGBColor(0xF3, 0xE5, 0xF5),
        "items": [
            "多模型交叉验证",
            "用例实际运行验证",
            "基于公司数据微调专属模型",
            "全链路自动化",
            "AI 主动发现漏测风险",
        ]
    },
]

card_w = Inches(3.7)
card_h = Inches(4.5)
gap = Inches(0.4)
start_x = Inches(0.6)
y = Inches(2.0)

for i, phase in enumerate(phases):
    x = start_x + i * (card_w + gap)

    card = add_rounded_rect(slide, x, y, card_w, card_h, phase["bg"])
    card.line.fill.solid()
    card.line.fill.fore_color.rgb = phase["color"]

    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.3), Inches(0.4),
                 phase["title"], font_size=17, bold=True, color=phase["color"])
    add_text_box(slide, x + Inches(0.15), y + Inches(0.5), card_w - Inches(0.3), Inches(0.3),
                 phase["subtitle"], font_size=13, color=MEDIUM_GRAY)

    tf = add_text_box(slide, x + Inches(0.15), y + Inches(0.9), card_w - Inches(0.3), card_h - Inches(1.1),
                      "", font_size=14, color=DARK)
    for item in phase["items"]:
        first_or_add(tf, f"✅  {item}" if i == 0 else f"•  {item}",
                     font_size=14, color=DARK, space_before=Pt(6))

# 箭头
add_text_box(slide, Inches(4.5), Inches(4.0), Inches(0.5), Inches(0.4),
             "→", font_size=30, bold=True, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.9), Inches(4.0), Inches(0.5), Inches(0.4),
             "→", font_size=30, bold=True, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5),
             "核心原则：先让 AI 帮忙分担 80% 的重复工作，再逐步提升 AI 用例的准确性和深度",
             font_size=16, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 11: 结束页
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, PRIMARY)
add_rect(slide, Inches(0.15), Inches(2.8), Inches(8), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.0),
             "Thanks", font_size=52, bold=True, color=DARK, align=PP_ALIGN.LEFT)

add_text_box(slide, Inches(0.8), Inches(3.1), Inches(11), Inches(0.6),
             "AiProtect — 让 AI 成为每位测试工程师的智能助手",
             font_size=22, color=MEDIUM_GRAY)

# STAR 总结
add_text_box(slide, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
             "STAR 总结", font_size=24, bold=True, color=PRIMARY)

star_summary = [
    ("S 场景", "测试用例编写耗时长、覆盖率难保证、格式不统一", PRIMARY),
    ("T 任务", "用 AI 自动生成高质量测试用例，释放测试人力", SECONDARY),
    ("A 行动", "3 步操作 + LangGraph 工作流 + 知识增强 + 双重验证", ORANGE),
    ("R 结果", "节省 80%+ 用例编写时间，覆盖率系统性保证", TEAL),
]

y_pos = Inches(4.9)
for label, text, color in star_summary:
    label_box = add_rounded_rect(slide, Inches(0.8), y_pos, Inches(1.2), Inches(0.35), color)
    label_box.line.fill.background()
    set_text(label_box, label, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    label_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, Inches(2.2), y_pos + Inches(0.02), Inches(9), Inches(0.35),
                 text, font_size=15, color=DARK)
    y_pos += Inches(0.45)

add_text_box(slide, Inches(3), Inches(6.8), Inches(7), Inches(0.5),
             "Q & A", font_size=36, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


# 保存
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "AiProtect_周会分享.pptx")
prs.save(output_path)
print(f"PPT 已生成: {os.path.abspath(output_path)}")
